"""
Data Ingestion module for RAG system
Module chuẩn bị và tải dữ liệu vào vector database
"""

import os
import asyncio
from pathlib import Path
from typing import List, Dict, Any, Optional, Union
from datetime import datetime
import mimetypes
from concurrent.futures import ThreadPoolExecutor

# Document processing
import PyPDF2
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation
import pandas as pd
from bs4 import BeautifulSoup
import markdown

# LangChain document loaders and splitters
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import (
    PyPDFLoader, Docx2txtLoader, UnstructuredExcelLoader,
    UnstructuredPowerPointLoader, TextLoader
)

from config import get_config
from utils import (
    setup_logging, clean_text, extract_metadata_from_filename,
    generate_document_id, chunk_text_with_overlap, validate_file_type,
    safe_filename, format_file_size, create_directory_if_not_exists,
    get_supported_file_types, is_text_meaningful
)
from embedding import EmbeddingManager
from database import QdrantManager

class DocumentProcessor:
    """
    Xử lý các loại documents khác nhau
    """
    
    def __init__(self):
        self.config = get_config()
        self.logger = setup_logging(self.config.app.log_level, self.config.app.log_file)
        self.supported_types = get_supported_file_types()
        
    def process_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        """Xử lý file PDF"""
        try:
            loader = PyPDFLoader(file_path)
            pages = loader.load()
            
            documents = []
            for i, page in enumerate(pages):
                content = clean_text(page.page_content)
                if is_text_meaningful(content):
                    documents.append({
                        'content': content,
                        'page_number': i + 1,
                        'source_type': 'pdf'
                    })
            
            self.logger.debug(f"✅ Đã xử lý PDF: {len(documents)} pages từ {Path(file_path).name}")
            return documents
            
        except Exception as e:
            self.logger.error(f"❌ Lỗi xử lý PDF {file_path}: {str(e)}")
            return []
    
    def process_docx(self, file_path: str) -> List[Dict[str, Any]]:
        """Xử lý file DOCX"""
        try:
            loader = Docx2txtLoader(file_path)
            doc = loader.load()[0]
            
            content = clean_text(doc.page_content)
            if is_text_meaningful(content):
                return [{
                    'content': content,
                    'source_type': 'docx'
                }]
            return []
            
        except Exception as e:
            self.logger.error(f"❌ Lỗi xử lý DOCX {file_path}: {str(e)}")
            return []
    
    def process_excel(self, file_path: str) -> List[Dict[str, Any]]:
        """Xử lý file Excel"""
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            documents = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                # Convert sheet to text
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                        rows.append(row_text)
                
                if rows:
                    content = '\n'.join(rows)
                    content = clean_text(content)
                    
                    if is_text_meaningful(content):
                        documents.append({
                            'content': content,
                            'sheet_name': sheet_name,
                            'source_type': 'excel'
                        })
            
            self.logger.debug(f"✅ Đã xử lý Excel: {len(documents)} sheets từ {Path(file_path).name}")
            return documents
            
        except Exception as e:
            self.logger.error(f"❌ Lỗi xử lý Excel {file_path}: {str(e)}")
            return []
    
    def process_powerpoint(self, file_path: str) -> List[Dict[str, Any]]:
        """Xử lý file PowerPoint"""
        try:
            presentation = Presentation(file_path)
            documents = []
            
            for i, slide in enumerate(presentation.slides):
                text_content = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_content.append(shape.text)
                
                if text_content:
                    content = '\n'.join(text_content)
                    content = clean_text(content)
                    
                    if is_text_meaningful(content):
                        documents.append({
                            'content': content,
                            'slide_number': i + 1,
                            'source_type': 'powerpoint'
                        })
            
            self.logger.debug(f"✅ Đã xử lý PowerPoint: {len(documents)} slides từ {Path(file_path).name}")
            return documents
            
        except Exception as e:
            self.logger.error(f"❌ Lỗi xử lý PowerPoint {file_path}: {str(e)}")
            return []
    
    def process_text(self, file_path: str) -> List[Dict[str, Any]]:
        """Xử lý file text thuần"""
        try:
            loader = TextLoader(file_path, encoding='utf-8')
            doc = loader.load()[0]
            
            content = clean_text(doc.page_content)
            if is_text_meaningful(content):
                return [{
                    'content': content,
                    'source_type': 'text'
                }]
            return []
            
        except Exception as e:
            self.logger.error(f"❌ Lỗi xử lý text {file_path}: {str(e)}")
            return []
    
    def process_markdown(self, file_path: str) -> List[Dict[str, Any]]:
        """Xử lý file Markdown"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                md_content = f.read()
            
            # Convert markdown to HTML then to text
            html = markdown.markdown(md_content)
            soup = BeautifulSoup(html, 'html.parser')
            content = clean_text(soup.get_text())
            
            if is_text_meaningful(content):
                return [{
                    'content': content,
                    'source_type': 'markdown'
                }]
            return []
            
        except Exception as e:
            self.logger.error(f"❌ Lỗi xử lý Markdown {file_path}: {str(e)}")
            return []
    
    def process_file(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Xử lý file dựa trên extension
        
        Args:
            file_path: Đường dẫn file
            
        Returns:
            List of processed documents
        """
        file_path = Path(file_path)
        extension = file_path.suffix.lower()
        
        # Map extension to processor
        processors = {
            '.pdf': self.process_pdf,
            '.docx': self.process_docx,
            '.doc': self.process_docx,
            '.xlsx': self.process_excel,
            '.xls': self.process_excel,
            '.pptx': self.process_powerpoint,
            '.ppt': self.process_powerpoint,
            '.txt': self.process_text,
            '.md': self.process_markdown,
            '.py': self.process_text,
            '.js': self.process_text,
            '.html': self.process_text,
            '.css': self.process_text,
            '.json': self.process_text,
            '.xml': self.process_text,
            '.yaml': self.process_text,
            '.yml': self.process_text
        }
        
        processor = processors.get(extension)
        if not processor:
            self.logger.warning(f"⚠️ Không hỗ trợ file type: {extension}")
            return []
        
        return processor(str(file_path))

class DataIngestionPipeline:
    """
    Pipeline chính cho data ingestion
    """
    
    def __init__(
        self,
        embedding_manager: EmbeddingManager,
        qdrant_manager: QdrantManager
    ):
        """
        Khởi tạo DataIngestionPipeline
        
        Args:
            embedding_manager: EmbeddingManager instance
            qdrant_manager: QdrantManager instance
        """
        self.config = get_config()
        self.logger = setup_logging(self.config.app.log_level, self.config.app.log_file)
        
        self.embedding_manager = embedding_manager
        self.qdrant_manager = qdrant_manager
        self.document_processor = DocumentProcessor()
        
        # Text splitter
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.config.rag.chunk_size,
            chunk_overlap=self.config.rag.chunk_overlap,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
        
        # Create directories
        create_directory_if_not_exists(self.config.app.data_dir)
        create_directory_if_not_exists(self.config.app.upload_dir)
        
        self.logger.info("✅ DataIngestionPipeline đã khởi tạo")
    
    def process_single_file(
        self, 
        file_path: str,
        metadata: Optional[Dict[str, Any]] = None
    ) -> List[Dict[str, Any]]:
        """
        Xử lý một file duy nhất
        
        Args:
            file_path: Đường dẫn file
            metadata: Metadata bổ sung (optional)
            
        Returns:
            List of processed documents with embeddings
        """
        try:
            file_path = Path(file_path)
            
            # Kiểm tra file tồn tại
            if not file_path.exists():
                self.logger.error(f"❌ File không tồn tại: {file_path}")
                return []
            
            # Kiểm tra file type
            supported_extensions = []
            for category, exts in get_supported_file_types().items():
                supported_extensions.extend(exts)
            
            if not validate_file_type(str(file_path), supported_extensions):
                self.logger.warning(f"⚠️ File type không được hỗ trợ: {file_path.suffix}")
                return []
            
            # Extract file metadata
            file_metadata = extract_metadata_from_filename(str(file_path))
            if metadata:
                file_metadata.update(metadata)
            
            # Process document
            raw_documents = self.document_processor.process_file(str(file_path))
            if not raw_documents:
                self.logger.warning(f"⚠️ Không thể xử lý file: {file_path.name}")
                return []
            
            # Process each document part
            processed_documents = []
            
            for doc_part in raw_documents:
                # Split text into chunks
                chunks = self.text_splitter.split_text(doc_part['content'])
                
                for i, chunk_text in enumerate(chunks):
                    if not is_text_meaningful(chunk_text):
                        continue
                    
                    # Create document
                    doc = {
                        'id': generate_document_id(chunk_text, file_metadata),
                        'content': chunk_text,
                        'filename': file_metadata['filename'],
                        'source': str(file_path),
                        'chunk_id': i,
                        'total_chunks': len(chunks),
                        'document_type': file_metadata['document_type'],
                        'created_at': datetime.now().isoformat(),
                        'file_size': file_metadata['size'],
                        'metadata': {
                            **file_metadata,
                            **doc_part,
                            'chunk_index': i
                        }
                    }
                    
                    # Generate embedding
                    embedding = self.embedding_manager.embed_text(chunk_text)
                    doc['embedding'] = embedding
                    
                    processed_documents.append(doc)
            
            self.logger.info(f"✅ Đã xử lý file {file_path.name}: {len(processed_documents)} chunks")
            return processed_documents
            
        except Exception as e:
            self.logger.error(f"❌ Lỗi xử lý file {file_path}: {str(e)}")
            return []
    
    def process_directory(
        self,
        directory_path: str,
        recursive: bool = True,
        file_patterns: Optional[List[str]] = None
    ) -> List[Dict[str, Any]]:
        """
        Xử lý tất cả files trong directory
        
        Args:
            directory_path: Đường dẫn thư mục
            recursive: Có xử lý đệ quy không
            file_patterns: Patterns để lọc files (optional)
            
        Returns:
            List of all processed documents
        """
        try:
            directory_path = Path(directory_path)
            
            if not directory_path.exists() or not directory_path.is_dir():
                self.logger.error(f"❌ Directory không tồn tại: {directory_path}")
                return []
            
            # Get all files
            if recursive:
                files = list(directory_path.rglob("*"))
            else:
                files = list(directory_path.glob("*"))
            
            # Filter files
            supported_extensions = []
            for category, exts in get_supported_file_types().items():
                supported_extensions.extend(exts)
            
            valid_files = [
                f for f in files 
                if f.is_file() and validate_file_type(str(f), supported_extensions)
            ]
            
            if file_patterns:
                filtered_files = []
                for pattern in file_patterns:
                    filtered_files.extend(directory_path.glob(pattern))
                valid_files = [f for f in valid_files if f in filtered_files]
            
            self.logger.info(f"📁 Tìm thấy {len(valid_files)} files để xử lý trong {directory_path}")
            
            # Process files
            all_documents = []
            
            for file_path in valid_files:
                documents = self.process_single_file(str(file_path))
                all_documents.extend(documents)
            
            self.logger.info(f"✅ Đã xử lý {len(valid_files)} files: {len(all_documents)} documents")
            return all_documents
            
        except Exception as e:
            self.logger.error(f"❌ Lỗi xử lý directory {directory_path}: {str(e)}")
            return []
    
    def ingest_documents(
        self,
        documents: List[Dict[str, Any]],
        batch_size: int = 100
    ) -> bool:
        """
        Ingest documents vào Qdrant
        
        Args:
            documents: List of processed documents
            batch_size: Batch size cho ingestion
            
        Returns:
            True nếu thành công
        """
        try:
            if not documents:
                self.logger.warning("⚠️ Không có documents để ingest")
                return True
            
            # Ensure collection exists
            vector_size = self.embedding_manager.get_vector_dimension()
            self.qdrant_manager.create_collection(
                vector_size=vector_size,
                distance_metric=self.config.qdrant.distance_metric
            )
            
            # Insert documents
            success = self.qdrant_manager.insert_documents(documents, batch_size)
            
            if success:
                self.logger.info(f"✅ Đã ingest {len(documents)} documents vào Qdrant")
            else:
                self.logger.error("❌ Lỗi khi ingest documents")
            
            return success
            
        except Exception as e:
            self.logger.error(f"❌ Lỗi ingest documents: {str(e)}")
            return False
    
    def ingest_file(self, file_path: str, metadata: Optional[Dict[str, Any]] = None) -> bool:
        """
        Ingest một file hoàn chỉnh (process + ingest)
        
        Args:
            file_path: Đường dẫn file
            metadata: Metadata bổ sung
            
        Returns:
            True nếu thành công
        """
        documents = self.process_single_file(file_path, metadata)
        if documents:
            return self.ingest_documents(documents)
        return False
    
    def ingest_directory(
        self,
        directory_path: str,
        recursive: bool = True,
        batch_size: int = 100
    ) -> bool:
        """
        Ingest toàn bộ directory (process + ingest)
        
        Args:
            directory_path: Đường dẫn thư mục
            recursive: Có xử lý đệ quy không
            batch_size: Batch size cho ingestion
            
        Returns:
            True nếu thành công
        """
        documents = self.process_directory(directory_path, recursive)
        if documents:
            return self.ingest_documents(documents, batch_size)
        return False
    
    def update_document(
        self,
        document_id: str,
        new_content: str,
        metadata: Optional[Dict[str, Any]] = None
    ) -> bool:
        """
        Cập nhật một document đã tồn tại
        
        Args:
            document_id: ID của document
            new_content: Nội dung mới
            metadata: Metadata mới
            
        Returns:
            True nếu thành công
        """
        try:
            # Generate new embedding
            embedding = self.embedding_manager.embed_text(new_content)
            
            # Prepare updated document
            updated_doc = {
                'id': document_id,
                'content': new_content,
                'embedding': embedding,
                'updated_at': datetime.now().isoformat()
            }
            
            if metadata:
                updated_doc['metadata'] = metadata
            
            # Update in Qdrant
            success = self.qdrant_manager.insert_documents([updated_doc])
            
            if success:
                self.logger.info(f"✅ Đã cập nhật document: {document_id}")
            else:
                self.logger.error(f"❌ Lỗi cập nhật document: {document_id}")
            
            return success
            
        except Exception as e:
            self.logger.error(f"❌ Lỗi cập nhật document {document_id}: {str(e)}")
            return False
    
    def delete_documents_by_source(self, source_path: str) -> bool:
        """
        Xóa tất cả documents từ một source
        
        Args:
            source_path: Đường dẫn source
            
        Returns:
            True nếu thành công
        """
        try:
            filters = {'source': source_path}
            success = self.qdrant_manager.delete_documents(filters=filters)
            
            if success:
                self.logger.info(f"✅ Đã xóa documents từ source: {source_path}")
            else:
                self.logger.error(f"❌ Lỗi xóa documents từ source: {source_path}")
            
            return success
            
        except Exception as e:
            self.logger.error(f"❌ Lỗi xóa documents: {str(e)}")
            return False
    
    def get_ingestion_stats(self) -> Dict[str, Any]:
        """
        Lấy thống kê về ingestion
        
        Returns:
            Dictionary chứa stats
        """
        try:
            collection_info = self.qdrant_manager.get_collection_info()
            
            stats = {
                'total_documents': collection_info.get('points_count', 0),
                'collection_name': self.qdrant_manager.collection_name,
                'vector_dimension': collection_info.get('vector_size', 0),
                'distance_metric': collection_info.get('distance_metric', 'unknown'),
                'collection_status': collection_info.get('status', 'unknown'),
                'embedding_cache_size': self.embedding_manager.get_cache_size()
            }
            
            return stats
            
        except Exception as e:
            self.logger.error(f"❌ Lỗi lấy stats: {str(e)}")
            return {}

# Factory function
def create_ingestion_pipeline(
    embedding_manager: EmbeddingManager,
    qdrant_manager: QdrantManager
) -> DataIngestionPipeline:
    """
    Factory function để tạo DataIngestionPipeline
    """
    return DataIngestionPipeline(embedding_manager, qdrant_manager)
